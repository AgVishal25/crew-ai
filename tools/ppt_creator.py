from langchain.tools import tool
import os
from pptx import Presentation




class PptCreator:
    def __init__(self):
        pass
    
    @tool("Create a powerpoint ppt ")
    def create_ppt(self,slides_dict:dict):
        """
        Accepts a json/dict input as a parameter.
        Creates a PowerPoint presentation from a dictionary of slides.
        
        Parameters:
        slides_dict (dict): Dictionary where each key is 'slide_1', 'slide_2', etc., and each value is a dictionary with:
                            - "title": Title for the slide.
                            - "content": Content text for the slide.
        """
        file = "ppt_file.pptx"
        prs = Presentation()
        for slide_key in sorted(slides_dict.keys(), key=lambda x: int(x.split('_')[1])):
            slide_data = slides_dict[slide_key]
            
            # Choose a slide layout with a title and content (typically layout index 1)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set the title and content for the slide.
            slide.shapes.title.text = slide_data.get("title", "No Title")
            slide.placeholders[1].text = slide_data.get("content", "No Content")
    
        prs.save(file)
        print(f"Presentation saved as {file}")
        return True